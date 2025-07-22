import os
import re
import webbrowser
from pathlib import Path
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, simpledialog
from collections import defaultdict
import io
import sys
import base64
import threading
import csv

# --- PyInstallerでのexe化に伴うエラー対策 ---
# scipy(numpy.f2py)がインポート時に標準出力/エラー出力を要求することがあり、
# GUIアプリではこれらがNoneの場合があるため、エラー回避のために設定します。
if sys.stdout is None:
    sys.stdout = open(os.devnull, 'w')
if sys.stderr is None:
    sys.stderr = open(os.devnull, 'w')
# -----------------------------------------

# --- 必要なライブラリの確認 ---
# 初回実行前にターミナルで以下のコマンドを実行してください:
# pip install pandas matplotlib numpy scipy openpyxl ttkbootstrap tkinterdnd2 Pillow pywin32 python-pptx google-generativeai
try:
    import pandas as pd
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg, NavigationToolbar2Tk
    import numpy as np
    from scipy.signal import find_peaks
    import ttkbootstrap as tb
    from ttkbootstrap.constants import *
    from tkinterdnd2 import DND_FILES, TkinterDnD
    from PIL import Image
    import pptx
    from pptx.util import Inches
    import google.generativeai as genai
except ImportError as e:
    # 必要なライブラリがない場合はエラーメッセージを表示して終了
    root = tk.Tk()
    root.withdraw()
    messagebox.showerror("ライブラリ不足", f"必要なライブラリが見つかりません: {e}\n\nターミナルで以下のコマンドを実行してください:\npip install pandas matplotlib numpy scipy openpyxl ttkbootstrap tkinterdnd2 Pillow pywin32 python-pptx google-generativeai")
    exit()

# =============================================================================
# === Gemini設定ダイアログクラス ===
# =============================================================================
class GeminiSettingsDialog(tb.Toplevel):
    """Gemini APIキーとモデル名を入力するためのダイアログ"""
    def __init__(self, parent):
        super().__init__(parent)
        self.title("Gemini 設定")
        self.parent = parent
        self.geometry("450x150")

        self.api_key_var = tk.StringVar(value=parent.gemini_api_key)
        self.model_name_var = tk.StringVar(value=parent.gemini_model_name)

        frame = tb.Frame(self, padding=10)
        frame.pack(expand=True, fill="both")

        tb.Label(frame, text="APIキー:").grid(row=0, column=0, padx=5, pady=5, sticky="w")
        tb.Entry(frame, textvariable=self.api_key_var, width=45).grid(row=0, column=1, padx=5, pady=5)

        tb.Label(frame, text="モデル名:").grid(row=1, column=0, padx=5, pady=5, sticky="w")
        tb.Entry(frame, textvariable=self.model_name_var, width=45).grid(row=1, column=1, padx=5, pady=5)
        
        tb.Button(frame, text="保存", command=self.save_settings, bootstyle=SUCCESS).grid(row=2, column=1, padx=5, pady=10, sticky="e")
        
        self.transient(parent)
        self.grab_set()
        self.wait_window(self)

    def save_settings(self):
        self.parent.gemini_api_key = self.api_key_var.get()
        self.parent.gemini_model_name = self.model_name_var.get()
        self.parent._save_gemini_config()
        self.destroy()

# =============================================================================
# === Geminiチャットウィンドウクラス ===
# =============================================================================
class GeminiChatWindow(tb.Toplevel):
    def __init__(self, parent, api_key, model_name, image):
        super().__init__(parent)
        self.title("Gemini とチャット")
        self.geometry("700x500")

        self.api_key = api_key
        self.model_name = model_name
        self.image = image
        self.chat = None

        # UIの作成
        main_frame = tb.Frame(self, padding=10)
        main_frame.pack(fill="both", expand=True)
        main_frame.rowconfigure(0, weight=1)
        main_frame.columnconfigure(0, weight=1)

        self.chat_history = tk.Text(main_frame, wrap="word", state="disabled", font=("Meiryo UI", 10))
        self.chat_history.grid(row=0, column=0, columnspan=2, sticky="nsew")
        
        vsb = tb.Scrollbar(main_frame, orient="vertical", command=self.chat_history.yview)
        vsb.grid(row=0, column=2, sticky="ns")
        self.chat_history.config(yscrollcommand=vsb.set)

        self.user_input = tb.Entry(main_frame, font=("Meiryo UI", 10))
        self.user_input.grid(row=1, column=0, padx=(0, 5), pady=5, sticky="ew")
        self.user_input.bind("<Return>", self.send_message)

        self.send_button = tb.Button(main_frame, text="送信", command=self.send_message, bootstyle=SUCCESS)
        self.send_button.grid(row=1, column=1, pady=5, sticky="ew")

        self.after(100, self.start_chat) # ウィンドウ表示後にチャットを開始

    def _append_message(self, sender, message):
        self.chat_history.config(state="normal")
        if sender == "You":
            self.chat_history.insert("end", f"あなた: {message}\n\n", "user")
        else: # Gemini or System
            self.chat_history.insert("end", f"{sender}: {message}\n\n", "gemini")
        self.chat_history.config(state="disabled")
        self.chat_history.see("end")
        self.chat_history.tag_config("user", foreground="cyan")
        self.chat_history.tag_config("gemini", foreground="white")

    def start_chat(self):
        try:
            genai.configure(api_key=self.api_key)
            model = genai.GenerativeModel(self.model_name)
            self.chat = model.start_chat()
            
            initial_prompt = "あなたデータ解析のアシスタントです。以下のグラフを見て、その主な特徴、傾向、注目すべき点について、専門家でなくても理解しやすいように日本語で簡潔に説明してください。その後、ユーザーからの質問に答える準備をしてください。"
            self._append_message("System", "グラフを送信中...")
            self.send_message(initial_prompt=initial_prompt)

        except Exception as e:
            messagebox.showerror("APIエラー", f"チャットの開始に失敗しました:\n{e}", parent=self)
            self.destroy()

    def send_message(self, event=None, initial_prompt=None):
        if initial_prompt:
            user_message = initial_prompt
            content = [user_message, self.image]
        else:
            user_message = self.user_input.get()
            if not user_message.strip():
                return
            self._append_message("You", user_message)
            self.user_input.delete(0, "end")
            content = user_message
        
        self.send_button.config(state="disabled")
        self.user_input.config(state="disabled")

        thread = threading.Thread(target=self._stream_response, args=(content,))
        thread.start()

    def _stream_response(self, content):
        try:
            self.after(0, self._start_gemini_message)
            response = self.chat.send_message(content, stream=True)
            for chunk in response:
                self.after(0, self._append_chunk, chunk.text)
        except Exception as e:
            self.after(0, self._append_chunk, f"\nエラーが発生しました: {e}")
        finally:
            self.after(0, self._finalize_response)

    def _start_gemini_message(self):
        self.chat_history.config(state="normal")
        last_line_start = self.chat_history.index("end-2l")
        self.chat_history.delete(last_line_start, "end") # "送信中..." or "考え中..."を削除
        self.chat_history.insert("end", "Gemini: ", "gemini")
        self.chat_history.config(state="disabled")

    def _append_chunk(self, text):
        self.chat_history.config(state="normal")
        self.chat_history.insert("end", text)
        self.chat_history.see("end")
        self.chat_history.config(state="disabled")

    def _finalize_response(self):
        self.chat_history.config(state="normal")
        self.chat_history.insert("end", "\n\n")
        self.chat_history.config(state="disabled")
        self.chat_history.see("end")
        self.send_button.config(state="normal")
        self.user_input.config(state="normal")
        self.user_input.focus_set()

# =============================================================================
# === ★ 新バージョン ★ データベース検索 & 閲覧ウィンドウクラス ===
# =============================================================================
class PeakDatabaseSearchWindow(tb.Toplevel):
    """Raman/XRDのExcelデータベース検索と結果表示、内容閲覧を行うウィンドウ"""
    def __init__(self, parent, method, db_path):
        super().__init__(parent)
        self.parent = parent
        self.method = method
        self.db_file_path = db_path
        self.database = []

        self.x_col_name = '2θ (degree)' if self.method == 'xrd' else 'Raman Shift (cm-1)'
        self.required_cols = ['Label', self.x_col_name, 'Intensity', 'Reference']

        self.title(f"{method.upper()} データベース検索 & 閲覧")
        self.geometry("800x700")

        main_frame = tb.Frame(self, padding=10)
        main_frame.pack(fill="both", expand=True)
        
        pane = tb.PanedWindow(main_frame, orient=VERTICAL)
        pane.pack(fill=BOTH, expand=True)

        search_pane = tb.Frame(pane, padding=5)
        search_pane.rowconfigure(1, weight=1)
        search_pane.columnconfigure(0, weight=1)
        pane.add(search_pane, weight=2)

        settings_frame = tb.LabelFrame(search_pane, text="検索設定", padding=10)
        settings_frame.grid(row=0, column=0, sticky="ew", pady=(0, 10))
        
        self.tolerance_var = tk.DoubleVar(value=1.0)
        self.prominence_var = tk.DoubleVar(value=0.1)

        tb.Label(settings_frame, text="ピーク許容誤差 (±):").pack(side="left", padx=5)
        tb.Entry(settings_frame, textvariable=self.tolerance_var, width=8).pack(side="left", padx=5)
        tb.Label(settings_frame, text="ピーク検出感度 (0-1):").pack(side="left", padx=5)
        tb.Entry(settings_frame, textvariable=self.prominence_var, width=8).pack(side="left", padx=5)
        tb.Button(settings_frame, text="検索実行", command=self._search_database, bootstyle=INFO).pack(side="left", padx=10)
        tb.Button(settings_frame, text="DBファイル編集", command=self._edit_database_file, bootstyle=(SECONDARY, OUTLINE)).pack(side="right", padx=5)

        results_frame = tb.LabelFrame(search_pane, text="検索結果 (複数選択可)", padding=10)
        results_frame.grid(row=1, column=0, sticky="nsew")
        results_frame.rowconfigure(0, weight=1)
        results_frame.columnconfigure(0, weight=1)

        cols = ("material", "matched_peaks", "total_peaks")
        self.results_tree = tb.Treeview(results_frame, columns=cols, show="headings", selectmode="extended")
        self.results_tree.heading("material", text="候補材料 (シート名)")
        self.results_tree.heading("matched_peaks", text="一致ピーク数")
        self.results_tree.heading("total_peaks", text="DBピーク数")
        self.results_tree.column("material", width=250)
        self.results_tree.grid(row=0, column=0, sticky="nsew")
        vsb_res = tb.Scrollbar(results_frame, orient="vertical", command=self.results_tree.yview)
        vsb_res.grid(row=0, column=1, sticky="ns")
        self.results_tree.config(yscrollcommand=vsb_res.set)
        
        results_action_frame = tb.Frame(results_frame)
        results_action_frame.grid(row=2, column=0, columnspan=2, pady=(10, 0), sticky="e")
        tb.Button(results_action_frame, text="グラフのRefをクリア", command=self._clear_references_on_main_graph, bootstyle=(WARNING, OUTLINE)).pack(side="left", padx=5)
        tb.Button(results_action_frame, text="選択材料をグラフに描画", command=self._plot_selection_on_main_graph, bootstyle=SUCCESS).pack(side="left")

        viewer_pane = tb.Frame(pane, padding=5)
        viewer_pane.rowconfigure(1, weight=1)
        viewer_pane.columnconfigure(0, weight=1)
        pane.add(viewer_pane, weight=3)

        db_viewer_frame = tb.LabelFrame(viewer_pane, text="データベース閲覧", padding=10)
        db_viewer_frame.grid(row=1, column=0, sticky="nsew")
        db_viewer_frame.rowconfigure(1, weight=1)
        db_viewer_frame.columnconfigure(0, weight=1)

        db_controls_frame = tb.Frame(db_viewer_frame)
        db_controls_frame.grid(row=0, column=0, sticky="ew", pady=(0, 5))
        tb.Label(db_controls_frame, text="材料を選択:").pack(side="left", padx=(0, 5))
        self.material_var = tk.StringVar()
        self.material_combo = tb.Combobox(db_controls_frame, textvariable=self.material_var, state="readonly", width=30)
        self.material_combo.pack(side="left")
        self.material_combo.bind("<<ComboboxSelected>>", self._display_material_data)

        db_cols = ("Label", "X", "Y", "Reference")
        self.db_data_tree = tb.Treeview(db_viewer_frame, columns=db_cols, show="headings")
        self.db_data_tree.heading("Label", text="ラベル")
        self.db_data_tree.heading("X", text=self.x_col_name)
        self.db_data_tree.heading("Y", text="相対強度")
        self.db_data_tree.heading("Reference", text="引用")
        self.db_data_tree.column("Label", width=80, anchor='center')
        self.db_data_tree.column("X", width=120, anchor='e')
        self.db_data_tree.column("Y", width=80, anchor='e')
        self.db_data_tree.column("Reference", width=200)
        self.db_data_tree.grid(row=1, column=0, sticky="nsew")
        vsb_db = tb.Scrollbar(db_viewer_frame, orient="vertical", command=self.db_data_tree.yview)
        vsb_db.grid(row=1, column=1, sticky="ns")
        self.db_data_tree.config(yscrollcommand=vsb_db.set)

        self._load_database()

    def _load_database(self):
        self.database.clear()
        material_names = []
        try:
            if not self.db_file_path.exists():
                messagebox.showwarning("DBなし", f"データベースファイルが見つかりません:\n{self.db_file_path}", parent=self)
                return
            db_sheets = pd.read_excel(self.db_file_path, sheet_name=None)
            for sheet_name, df in db_sheets.items():
                if all(col in df.columns for col in self.required_cols):
                    self.database.append({"name": sheet_name, "data": df})
                    material_names.append(sheet_name)
                else:
                    print(f"警告: シート '{sheet_name}' は必須列({', '.join(self.required_cols)})を持っていないため、スキップします。")
            
            self.material_combo['values'] = material_names
            if material_names:
                self.material_var.set(material_names[0])
                self._display_material_data()

        except Exception as e:
            messagebox.showerror("データベース読込エラー", f"データベースの読み込みに失敗しました:\n{e}", parent=self)

    def _display_material_data(self, event=None):
        self.db_data_tree.delete(*self.db_data_tree.get_children())
        selected_name = self.material_var.get()
        material = next((m for m in self.database if m["name"] == selected_name), None)
        if material:
            df = material["data"]
            for _, row in df.iterrows():
                self.db_data_tree.insert("", "end", values=(
                    row["Label"], 
                    f"{row[self.x_col_name]:.2f}", 
                    f"{row['Intensity']:.1f}",
                    row["Reference"]
                ))

    def _edit_database_file(self):
        try:
            os.startfile(self.db_file_path)
        except Exception as e:
            messagebox.showerror("エラー", f"ファイルを開けませんでした:\n{e}", parent=self)

    def _find_peaks_in_plot(self):
        if not self.parent.fig or not self.parent.fig.axes:
            messagebox.showwarning("エラー", "グラフが描画されていません。", parent=self)
            return None
        ax = self.parent.fig.axes[0]
        if not ax.lines:
            messagebox.showwarning("エラー", "グラフにデータがありません。", parent=self)
            return None
        all_x, all_y = [], []
        for line in ax.lines:
            if not line.get_label().startswith('_') and not "Ref:" in line.get_label():
                 x, y = line.get_data()
                 all_x.extend(x)
                 all_y.extend(y)
        if not all_x: return None
        sorted_indices = np.argsort(all_x)
        x_data, y_data = np.array(all_x)[sorted_indices], np.array(all_y)[sorted_indices]
        prominence = self.prominence_var.get() * (np.max(y_data) - np.min(y_data))
        peak_indices, _ = find_peaks(y_data, prominence=prominence)
        return x_data[peak_indices]

    def _search_database(self):
        measured_peaks = self._find_peaks_in_plot()
        if measured_peaks is None: return
        self._load_database()
        self.results_tree.delete(*self.results_tree.get_children())
        tolerance = self.tolerance_var.get()
        results = []
        for material in self.database:
            match_count = 0
            for db_peak in material["data"][self.x_col_name]:
                if any(abs(measured_peak - db_peak) <= tolerance for measured_peak in measured_peaks):
                    match_count += 1
            if match_count > 0:
                results.append({"name": material["name"], "matched": match_count, "total": len(material["data"][self.x_col_name])})
        results.sort(key=lambda x: x["matched"], reverse=True)
        if not results:
            messagebox.showinfo("結果なし", "一致する材料がデータベースに見つかりませんでした。", parent=self)
            return
        for res in results:
            self.results_tree.insert("", "end", values=(res["name"], res["matched"], res["total"]))

    def _plot_selection_on_main_graph(self):
        selected_items = self.results_tree.selection()
        if not selected_items:
            messagebox.showwarning("未選択", "グラフに描画する材料を1つ以上選択してください。", parent=self)
            return
        
        self.parent._clear_reference_artists()

        for item_id in selected_items:
            selected_material_name = self.results_tree.item(item_id, "values")[0]
            material_data = next((m for m in self.database if m["name"] == selected_material_name), None)
            if material_data:
                self.parent._plot_reference_spectrum(material_data["data"], material_data["name"], add_mode=True)
        
        self.parent._update_legend()

    def _clear_references_on_main_graph(self):
        self.parent._clear_reference_artists()


# =============================================================================
# === データ解析アプリケーションクラス ===
# =============================================================================
class DataAnalyzerApp(TkinterDnD.Tk):
    def __init__(self):
        super().__init__()
        self.withdraw()
        tb.Style(theme="darkly")
        self.title("汎用データ解析ツール")
        self.geometry("1300x800")
        self.minsize(900, 600)
        self.canvas_widget, self.toolbar, self.fig = None, None, None
        self.plot_files, self.current_folder_path, self.temp_ppt_path = [], None, None
        self.gemini_api_key, self.gemini_model_name = "", "gemini-1.5-flash"
        self.config_file = Path("gemini_config.txt")
        self.reference_artists = []
        
        self.ref_colors = ['#ff7f0e', '#2ca02c', '#d62728', '#9467bd', '#8c564b', '#e377c2', '#7f7f7f', '#bcbd22', '#17becf']
        self.ref_color_index = 0
        
        self.xrd_db_file = Path("xrd_database.xlsx")
        self.raman_db_file = Path("raman_database.xlsx")
        self._create_database_files_if_not_exist()

        self.method_var = tk.StringVar(value='xrd')
        self.auto_detect_header_var = tk.BooleanVar(value=True)
        self.skip_rows_var = tk.IntVar(value=1)
        self.x_min_var, self.x_max_var = tk.StringVar(), tk.StringVar()
        self.y_min_var, self.y_max_var = tk.StringVar(), tk.StringVar()
        self.current_path_var = tk.StringVar(value="フォルダが選択されていません")
        self.log_scale_var = tk.BooleanVar(value=False)
        self.abs_current_var = tk.BooleanVar(value=False)
        self.individual_offset_op_var = tk.StringVar(value='+')
        self.individual_offset_val_var = tk.StringVar(value="0.0")
        self.x_axis_var = tk.StringVar(value='Voltage')
        self.y_axis_var = tk.StringVar(value='Current')
        self.area_var = tk.StringVar(value='1.0')
        self.actual_depth_var = tk.StringVar()
        self.ftir_original_y_unit_var = tk.StringVar(value='Absorbance')
        self.ftir_target_y_unit_var = tk.StringVar(value='Absorbance')

        self._load_gemini_config()
        self._create_main_widgets()
        self.protocol("WM_DELETE_WINDOW", self.destroy)
        self.deiconify()

    def _create_database_files_if_not_exist(self):
        """初回起動時にXRDとRamanのExcelデータベースファイルを個別に作成する"""
        if not self.xrd_db_file.exists():
            try:
                with pd.ExcelWriter(self.xrd_db_file, engine='openpyxl') as writer:
                    df_si = pd.DataFrame({
                        'Label': ['(111)', '(220)', '(311)'], 
                        '2θ (degree)': [28.4, 47.3, 56.1], 
                        'Intensity': [100, 60, 40],
                        'Reference': ['COD: 9011422', 'COD: 9011422', 'COD: 9011422']
                    })
                    df_si.to_excel(writer, sheet_name='Si', index=False)
                    df_tio2 = pd.DataFrame({
                        'Label': ['A(101)', 'R(110)', 'A(200)'], 
                        '2θ (degree)': [25.3, 27.4, 48.1], 
                        'Intensity': [100, 80, 50],
                        'Reference': ['JCPDS 21-1272', 'JCPDS 21-1276', 'JCPDS 21-1272']
                    })
                    df_tio2.to_excel(writer, sheet_name='TiO2 (Anatase-Rutile)', index=False)
            except Exception as e:
                messagebox.showerror("DB作成エラー", f"XRDデータベース作成に失敗:\n{e}", parent=self)
        
        if not self.raman_db_file.exists():
            try:
                with pd.ExcelWriter(self.raman_db_file, engine='openpyxl') as writer:
                    df_si = pd.DataFrame({
                        'Label': ['Si'], 
                        'Raman Shift (cm-1)': [520.7], 
                        'Intensity': [100],
                        'Reference': ['Parker et al. 1967']
                    })
                    df_si.to_excel(writer, sheet_name='Si', index=False)
                    df_mos2 = pd.DataFrame({
                        'Label': ['E2g', 'A1g'], 
                        'Raman Shift (cm-1)': [383, 408], 
                        'Intensity': [90, 100],
                        'Reference': ['Li et al. 2012']
                    })
                    df_mos2.to_excel(writer, sheet_name='MoS2', index=False)
            except Exception as e:
                messagebox.showerror("DB作成エラー", f"Ramanデータベース作成に失敗:\n{e}", parent=self)

    def _load_gemini_config(self):
        if self.config_file.exists():
            try:
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                    if len(lines) >= 2:
                        self.gemini_api_key = lines[0].strip()
                        self.gemini_model_name = lines[1].strip()
            except Exception: pass

    def _save_gemini_config(self):
        try:
            with open(self.config_file, 'w', encoding='utf-8') as f:
                f.write(f"{self.gemini_api_key}\n")
                f.write(f"{self.gemini_model_name}\n")
        except Exception as e:
            messagebox.showwarning("設定保存エラー", f"設定の保存に失敗しました:\n{e}", parent=self)

    def _create_main_widgets(self):
        top_controls = tb.Frame(self, padding=(10, 10, 10, 0))
        top_controls.pack(fill=tk.X)
        tb.Button(top_controls, text="データフォルダを開く", command=self._select_folder, bootstyle=INFO).pack(side=tk.LEFT, padx=5)
        tb.Label(top_controls, textvariable=self.current_path_var, relief="sunken", anchor="w").pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)
        tb.Label(top_controls, text="測定方法:").pack(side=tk.LEFT, padx=(10, 5))
        method_combo = tb.Combobox(top_controls, textvariable=self.method_var, values=['xrd', 'raman', 'iv', 'ftir', 'xps', 'zem3'], state='readonly')
        method_combo.pack(side=tk.LEFT, padx=5)
        method_combo.bind("<<ComboboxSelected>>", self._on_method_change)
        main_pane = tb.PanedWindow(self, orient=tk.HORIZONTAL)
        main_pane.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        left_frame = tb.Frame(main_pane, padding=5)
        main_pane.add(left_frame, weight=1)
        self._create_file_list_widgets(left_frame)
        self.right_frame = tb.Frame(main_pane, padding=5)
        main_pane.add(self.right_frame, weight=3)
        self._create_graph_area_widgets(self.right_frame)

    def _select_folder(self):
        folder = filedialog.askdirectory(title="解析したいデータフォルダを選択")
        if folder: self._change_directory(Path(folder))

    def _on_method_change(self, event=None):
        for widget in self.graph_controls_container.winfo_children(): widget.destroy()
        self._create_graph_controls(self.graph_controls_container)
        self._update_item_settings_ui()
        new_method = self.method_var.get()
        for item in self.plot_files:
            if new_method == 'zem3':
                if 'actual_depth_cm' not in item: item['actual_depth_cm'] = self._read_zem_depth(item['path'])
            else:
                if 'op' not in item: item['op'] = '+'
                if 'val' not in item: item['val'] = 0.0
        self.plot_list_tree.delete(*self.plot_list_tree.get_children())
        for item in self.plot_files:
            file_path = item['path']
            if new_method == 'zem3':
                depth = item['actual_depth_cm']
                self.plot_list_tree.insert("", "end", text=file_path.name, values=(f"深さ: {depth:.4g} cm",))
            else:
                op, val = item['op'], item['val']
                self.plot_list_tree.insert("", "end", text=file_path.name, values=(f"Y {op} {val:.4g}",))
        if self.plot_files: self._plot_from_list()

    def _create_file_list_widgets(self, parent):
        parent.rowconfigure(0, weight=1); parent.rowconfigure(2, weight=1); parent.columnconfigure(0, weight=1)
        browser_frame = tb.LabelFrame(parent, text="ファイルブラウザ (ファイル/フォルダをここにドロップ)", padding=5)
        browser_frame.grid(row=0, column=0, sticky="nsew"); browser_frame.rowconfigure(1, weight=1); browser_frame.columnconfigure(0, weight=1)
        tree_frame = tb.Frame(browser_frame); tree_frame.grid(row=1, column=0, sticky="nsew"); tree_frame.rowconfigure(0, weight=1); tree_frame.columnconfigure(0, weight=1)
        self.file_tree = tb.Treeview(tree_frame, columns=("type", "size", "modified"), selectmode="extended", bootstyle=INFO)
        self.file_tree.heading("#0", text="名前"); self.file_tree.heading("type", text="種類"); self.file_tree.heading("size", text="サイズ"); self.file_tree.heading("modified", text="更新日時")
        self.file_tree.column("#0", width=150, stretch=True); self.file_tree.column("type", width=50, anchor="center", stretch=False); self.file_tree.column("size", width=70, anchor="e", stretch=False); self.file_tree.column("modified", width=120, anchor="center", stretch=False)
        vsb = tb.Scrollbar(tree_frame, orient="vertical", command=self.file_tree.yview, bootstyle=INFO)
        self.file_tree.configure(yscrollcommand=vsb.set); self.file_tree.grid(row=0, column=0, sticky="nsew"); vsb.grid(row=0, column=1, sticky="ns"); self.file_tree.bind("<Double-1>", self._on_item_double_click)
        self.file_tree.drop_target_register(DND_FILES)
        self.file_tree.dnd_bind('<<Drop>>', self.on_drop)
        buttons_frame = tb.Frame(parent); buttons_frame.grid(row=1, column=0, pady=10)
        tb.Button(buttons_frame, text="▶ 描画リストに追加", command=self._add_to_plot_list, bootstyle=SUCCESS).pack(pady=2)
        tb.Button(buttons_frame, text="◀ 選択項目を削除", command=self._remove_from_plot_list, bootstyle=(DANGER, OUTLINE)).pack(pady=2)
        tb.Button(buttons_frame, text="✖ 全て削除", command=self._clear_plot_list, bootstyle=(WARNING, OUTLINE)).pack(pady=2)
        plot_list_frame = tb.LabelFrame(parent, text="描画リスト", padding=5); plot_list_frame.grid(row=2, column=0, sticky="nsew"); plot_list_frame.rowconfigure(0, weight=1); plot_list_frame.columnconfigure(0, weight=1)
        self.plot_list_tree = tb.Treeview(plot_list_frame, columns=("info",), selectmode="extended", bootstyle=PRIMARY)
        self.plot_list_tree.heading("info", text="情報"); self.plot_list_tree.column("info", width=120, anchor="center", stretch=False)
        self.plot_list_tree.heading("#0", text="ファイル名"); self.plot_list_tree.column("#0", stretch=True); self.plot_list_tree.grid(row=0, column=0, sticky="nsew")
        plot_vsb = tb.Scrollbar(plot_list_frame, orient="vertical", command=self.plot_list_tree.yview, bootstyle=PRIMARY)
        plot_vsb.grid(row=0, column=1, sticky="ns"); self.plot_list_tree.config(yscrollcommand=plot_vsb.set)
        self.plot_list_tree.bind("<<TreeviewSelect>>", self._on_plot_list_select)
        self.plot_list_tree.bind("<Delete>", lambda e: self._remove_from_plot_list())
        plot_order_frame = tb.Frame(plot_list_frame); plot_order_frame.grid(row=0, column=2, padx=5, sticky="ns")
        tb.Button(plot_order_frame, text="▲", command=lambda: self._move_in_plot_list(-1), bootstyle=(SECONDARY, OUTLINE)).pack(pady=2, fill=tk.X)
        tb.Button(plot_order_frame, text="▼", command=lambda: self._move_in_plot_list(1), bootstyle=(SECONDARY, OUTLINE)).pack(pady=2, fill=tk.X)
        self.item_settings_frame = tb.Frame(plot_list_frame); self.item_settings_frame.grid(row=1, column=0, columnspan=3, sticky="ew", pady=(5,0))
        self._update_item_settings_ui()

    def _create_graph_area_widgets(self, parent):
        parent.rowconfigure(0, weight=1); parent.columnconfigure(0, weight=1)
        self.plot_frame = tb.Frame(parent, relief="sunken", borderwidth=1); self.plot_frame.grid(row=0, column=0, sticky="nsew", pady=(0, 10))
        self.graph_controls_container = tb.LabelFrame(parent, text="グラフ設定", padding=10)
        self.graph_controls_container.grid(row=1, column=0, sticky="ew")
        self._create_graph_controls(self.graph_controls_container)

    def _create_graph_controls(self, parent):
        parent.columnconfigure(1, weight=1)
        method = self.method_var.get()
        action_buttons_frame = tb.Frame(parent)
        action_buttons_frame.grid(row=0, column=0, rowspan=2, sticky='ns', padx=(0, 15))
        tb.Button(action_buttons_frame, text="グラフ描画", command=self._plot_from_list, bootstyle=PRIMARY).pack(fill=tk.X, pady=2, ipady=2)
        if method == 'zem3':
            tb.Button(action_buttons_frame, text="データテーブル表示", command=self._show_zem3_table_popup, bootstyle=(SECONDARY, OUTLINE)).pack(fill=tk.X, pady=2, ipady=2)
        else:
            if method in ['xrd', 'raman']:
                tb.Button(action_buttons_frame, text="データベース検索", command=self._open_database_search_window, bootstyle=(INFO, OUTLINE)).pack(fill=tk.X, pady=2, ipady=2)
        tb.Button(action_buttons_frame, text="PPTへ出力", command=self._export_to_ppt, bootstyle=(SUCCESS, OUTLINE)).pack(fill=tk.X, pady=2, ipady=2)
        tb.Button(action_buttons_frame, text="Geminiとチャット", command=self._start_gemini_chat, bootstyle=(INFO, OUTLINE)).pack(fill=tk.X, pady=2, ipady=2)
        method_options_frame = tb.Frame(parent)
        method_options_frame.grid(row=0, column=1, sticky='ew')
        if method == 'iv':
            iv_controls_frame = tb.Frame(method_options_frame); iv_controls_frame.pack(side=tk.LEFT)
            tb.Label(iv_controls_frame, text="横軸(X):").grid(row=0, column=0, padx=(0, 2), pady=2); self.x_axis_combo = tb.Combobox(iv_controls_frame, textvariable=self.x_axis_var, values=['Voltage', 'Time'], state='readonly', width=10); self.x_axis_combo.grid(row=0, column=1, pady=2)
            tb.Label(iv_controls_frame, text="縦軸(Y):").grid(row=0, column=2, padx=(5, 2), pady=2); self.y_axis_combo = tb.Combobox(iv_controls_frame, textvariable=self.y_axis_var, values=['Current', 'Current Density', 'Resistance'], state='readonly', width=15); self.y_axis_combo.grid(row=0, column=3, pady=2); self.y_axis_combo.bind("<<ComboboxSelected>>", self._toggle_area_entry)
            self.area_label = tb.Label(iv_controls_frame, text="面積 (cm²):"); self.area_label.grid(row=0, column=4, padx=(5, 2), pady=2); self.area_entry = tb.Entry(iv_controls_frame, textvariable=self.area_var, width=8); self.area_entry.grid(row=0, column=5, pady=2); self._toggle_area_entry()
            tb.Checkbutton(iv_controls_frame, text="絶対値", variable=self.abs_current_var, bootstyle=(SQUARE, TOGGLE, INFO)).grid(row=0, column=6, padx=(10, 0))
        elif method == 'ftir':
            ftir_controls_frame = tb.Frame(method_options_frame); ftir_controls_frame.pack(side=tk.LEFT)
            tb.Label(ftir_controls_frame, text="元の単位:").grid(row=0, column=0, padx=(0, 2), pady=2)
            original_unit_combo = tb.Combobox(ftir_controls_frame, textvariable=self.ftir_original_y_unit_var, values=['Absorbance', 'Transmittance (%)', 'Reflectance (%)'], state='readonly', width=15); original_unit_combo.grid(row=0, column=1, pady=2)
            tb.Label(ftir_controls_frame, text="縦軸(Y):").grid(row=0, column=2, padx=(5, 2), pady=2)
            target_unit_combo = tb.Combobox(ftir_controls_frame, textvariable=self.ftir_target_y_unit_var, values=['Absorbance', 'Transmittance (%)', 'Reflectance (%)'], state='readonly', width=15); target_unit_combo.grid(row=0, column=3, pady=2)
        elif method == 'zem3':
            zem3_controls_frame = tb.Frame(method_options_frame); zem3_controls_frame.pack(side=tk.LEFT)
            zem3_cols = ['Temperature (°C)', 'Electrical conductivity (S/cm)', 'Seebeck coefficient (μV/K)', 'Power Factor (μW/cmK²)']
            tb.Label(zem3_controls_frame, text="横軸(X):").grid(row=0, column=0, padx=(0, 2), pady=2)
            self.zem3_x_axis_combo = tb.Combobox(zem3_controls_frame, textvariable=self.x_axis_var, values=zem3_cols, state='readonly', width=30); self.zem3_x_axis_combo.grid(row=0, column=1, pady=2); self.x_axis_var.set('Temperature (°C)')
            tb.Label(zem3_controls_frame, text="縦軸(Y):").grid(row=0, column=2, padx=(5, 2), pady=2)
            self.zem3_y_axis_combo = tb.Combobox(zem3_controls_frame, textvariable=self.y_axis_var, values=zem3_cols, state='readonly', width=30); self.zem3_y_axis_combo.grid(row=0, column=3, pady=2); self.y_axis_var.set('Electrical conductivity (S/cm)')
        if method not in ['zem3']:
            header_controls_frame = tb.Frame(method_options_frame); header_controls_frame.pack(side=tk.LEFT, padx=(10,0))
            auto_detect_cb = tb.Checkbutton(header_controls_frame, text="ヘッダーを自動検出", variable=self.auto_detect_header_var, bootstyle=(SQUARE, TOGGLE, INFO), command=self._toggle_skip_rows_spinbox); auto_detect_cb.pack(side=tk.LEFT, padx=(0, 5))
            tb.Label(header_controls_frame, text="データ開始行:").pack(side=tk.LEFT)
            self.skip_rows_spinbox = tb.Spinbox(header_controls_frame, from_=1, to=200, textvariable=self.skip_rows_var, width=5, bootstyle=PRIMARY); self.skip_rows_spinbox.pack(side=tk.LEFT, padx=5)
            self._toggle_skip_rows_spinbox()
        common_controls_frame = tb.Frame(parent); common_controls_frame.grid(row=1, column=1, sticky='ew', pady=(5,0))
        tb.Checkbutton(common_controls_frame, text="Logスケール", variable=self.log_scale_var, bootstyle=(SQUARE, TOGGLE, INFO)).pack(side=tk.LEFT, padx=(0, 15))
        for label, var in [("X min:", self.x_min_var), ("X max:", self.x_max_var), ("Y min:", self.y_min_var), ("Y max:", self.y_max_var)]:
            tb.Label(common_controls_frame, text=label).pack(side=tk.LEFT, padx=(10, 0))
            tb.Entry(common_controls_frame, textvariable=var, width=8).pack(side=tk.LEFT, padx=5)

    def _update_item_settings_ui(self, event=None):
        for widget in self.item_settings_frame.winfo_children(): widget.destroy()
        method = self.method_var.get()
        if method == 'zem3':
            tb.Label(self.item_settings_frame, text="選択項目の深さ (cm):").pack(side=tk.LEFT)
            self.depth_entry = tb.Entry(self.item_settings_frame, textvariable=self.actual_depth_var, width=10); self.depth_entry.pack(side=tk.LEFT, padx=5)
            tb.Button(self.item_settings_frame, text="設定", command=self._set_item_depth, bootstyle=(INFO, OUTLINE)).pack(side=tk.LEFT)
        else:
            tb.Label(self.item_settings_frame, text="選択項目のY軸オフセット:").pack(side=tk.LEFT)
            self.op_combo = tb.Combobox(self.item_settings_frame, textvariable=self.individual_offset_op_var, values=['+', '*'], width=2, state='readonly'); self.op_combo.pack(side=tk.LEFT, padx=(5, 2))
            self.offset_entry = tb.Entry(self.item_settings_frame, textvariable=self.individual_offset_val_var, width=8); self.offset_entry.pack(side=tk.LEFT, padx=2)
            tb.Button(self.item_settings_frame, text="設定", command=self._set_item_operation, bootstyle=(INFO, OUTLINE)).pack(side=tk.LEFT)

    def _toggle_area_entry(self, event=None):
        state = 'normal' if self.y_axis_var.get() == 'Current Density' else 'disabled'
        self.area_entry.config(state=state); self.area_label.config(state=state)

    def _toggle_skip_rows_spinbox(self):
        state = 'disabled' if self.auto_detect_header_var.get() else 'normal'
        if hasattr(self, 'skip_rows_spinbox'): self.skip_rows_spinbox.config(state=state)

    def _change_directory(self, target_path: Path):
        self.current_folder_path = target_path.resolve()
        self.current_path_var.set(str(self.current_folder_path))
        for item in self.file_tree.get_children(): self.file_tree.delete(item)
        if self.current_folder_path.parent != self.current_folder_path:
            self.file_tree.insert("", "end", text=".. (親フォルダへ)", values=("", "", ""), tags=('nav',))
        try:
            items = sorted(list(self.current_folder_path.iterdir()), key=lambda p: (p.is_file(), p.name.lower()))
            for path in items:
                try:
                    if path.is_dir(): self.file_tree.insert("", "end", text=path.name, values=("フォルダ", "", ""), tags=('folder',))
                    elif path.is_file():
                        size = f"{path.stat().st_size / 1024:.1f} KB"; modified = pd.to_datetime(path.stat().st_mtime, unit='s').strftime('%Y-%m-%d %H:%M')
                        self.file_tree.insert("", "end", text=path.name, values=("ファイル", size, modified), tags=('file',))
                except Exception: continue
        except Exception as e:
            messagebox.showerror("アクセスエラー", f"フォルダ '{self.current_folder_path}' の読み込みに失敗しました。\n詳細: {e}", parent=self)

    def _on_item_double_click(self, event):
        selected_id = self.file_tree.focus()
        if not selected_id: return
        item_text = self.file_tree.item(selected_id, "text")
        if item_text == ".. (親フォルダへ)":
            self._change_directory(self.current_folder_path.parent)
            return
        target_path = self.current_folder_path / item_text
        if target_path.is_dir(): self._change_directory(target_path)
        elif target_path.is_file(): self._add_single_file_to_plot_list(target_path)

    def on_drop(self, event):
        try:
            filepaths = self.tk.splitlist(event.data)
            for f in filepaths:
                p = Path(f)
                if p.is_file(): self._add_single_file_to_plot_list(p)
                elif p.is_dir():
                    for child in p.rglob('*'):
                        if child.is_file(): self._add_single_file_to_plot_list(child)
        except Exception as e:
            messagebox.showerror("ドロップエラー", f"ファイルの処理中にエラーが発生しました: {e}", parent=self)

    def _detect_data_start_row(self, file_path, consecutive_threshold=5):
        data_line_regex = re.compile(r'^\s*([-+]?\d*\.?\d+(?:[eE][-+]?\d+)?\s*[,;\s\t]\s*)+[-+]?\d*\.?\d+(?:[eE][-+]?\d+)?\s*$')
        try:
            lines = []
            try:
                with open(file_path, 'r', encoding='utf-8') as f: lines = f.readlines(4096)
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='cp932') as f: lines = f.readlines(4096)
            consecutive_count, potential_start_line = 0, -1
            for i, line in enumerate(lines):
                if data_line_regex.match(line.strip()):
                    if consecutive_count == 0: potential_start_line = i
                    consecutive_count += 1
                    if consecutive_count >= consecutive_threshold: return potential_start_line
                else: consecutive_count = 0
            if potential_start_line != -1: return potential_start_line
        except Exception: return 0
        return 0

    def _add_to_plot_list(self):
        for item_id in self.file_tree.selection():
            if self.file_tree.item(item_id, "values")[0] == "ファイル":
                file_path = self.current_folder_path / self.file_tree.item(item_id, "text")
                self._add_single_file_to_plot_list(file_path)

    def _add_single_file_to_plot_list(self, file_path: Path):
        if not any(d['path'] == file_path for d in self.plot_files):
            method = self.method_var.get()
            if method == 'zem3':
                depth = self._read_zem_depth(file_path)
                item = {'path': file_path, 'actual_depth_cm': depth}
                self.plot_list_tree.insert("", "end", text=file_path.name, values=(f"深さ: {depth:.4g} cm",))
            else:
                item = {'path': file_path, 'op': '+', 'val': 0.0}
                self.plot_list_tree.insert("", "end", text=file_path.name, values=(f"Y {item['op']} {item['val']:.4g}",))
            self.plot_files.append(item)

    def _read_zem_depth(self, file_path):
        try:
            with open(file_path, 'r', encoding='cp932') as f:
                match = re.search(r"Depth=([\d.E+-]+)", f.readline())
                return float(match.group(1)) * 100 if match else 0.0
        except Exception: return 0.0

    def _remove_from_plot_list(self):
        selected_iids = self.plot_list_tree.selection()
        if not selected_iids: return
        indices_to_remove = sorted([self.plot_list_tree.index(iid) for iid in selected_iids], reverse=True)
        for index in indices_to_remove: del self.plot_files[index]
        for iid in selected_iids: self.plot_list_tree.delete(iid)

    def _clear_plot_list(self):
        self.plot_files.clear()
        self.plot_list_tree.delete(*self.plot_list_tree.get_children())

    def _move_in_plot_list(self, direction):
        for iid in self.plot_list_tree.selection():
            index = self.plot_list_tree.index(iid)
            new_index = index + direction
            if 0 <= new_index < len(self.plot_files):
                self.plot_files.insert(new_index, self.plot_files.pop(index))
                self.plot_list_tree.move(iid, self.plot_list_tree.parent(iid), new_index)

    def _on_plot_list_select(self, event):
        if len(self.plot_list_tree.selection()) == 1:
            index = self.plot_list_tree.index(self.plot_list_tree.selection()[0])
            item_data = self.plot_files[index]
            method = self.method_var.get()
            if method == 'zem3': self.actual_depth_var.set(str(item_data.get('actual_depth_cm', '')))
            else:
                self.individual_offset_op_var.set(item_data.get('op', '+'))
                self.individual_offset_val_var.set(str(item_data.get('val', 0.0)))

    def _set_item_operation(self):
        selected_iids = self.plot_list_tree.selection()
        if not selected_iids: messagebox.showwarning("項目未選択", "操作を設定する項目を選択してください。", parent=self); return
        try: new_val = float(self.individual_offset_val_var.get())
        except ValueError: messagebox.showerror("入力エラー", "操作の値には数値を入力してください。", parent=self); return
        op = self.individual_offset_op_var.get()
        for iid in selected_iids:
            index = self.plot_list_tree.index(iid)
            self.plot_files[index]['op'], self.plot_files[index]['val'] = op, new_val
            self.plot_list_tree.item(iid, values=(f"Y {op} {new_val:.4g}",))

    def _set_item_depth(self):
        selected_iids = self.plot_list_tree.selection()
        if not selected_iids: messagebox.showwarning("項目未選択", "深さを設定する項目を選択してください。", parent=self); return
        try: new_depth = float(self.actual_depth_var.get())
        except ValueError: messagebox.showerror("入力エラー", "深さには数値を入力してください。", parent=self); return
        for iid in selected_iids:
            index = self.plot_list_tree.index(iid)
            self.plot_files[index]['actual_depth_cm'] = new_depth
            self.plot_list_tree.item(iid, values=(f"深さ: {new_depth:.4g} cm",))

    def _clear_reference_artists(self):
        """★ 改善版 ★ グラフ上のリファレンス関連描画オブジェクトをすべて削除し、凡例を更新する"""
        for artist in self.reference_artists:
            try: artist.remove()
            except (ValueError, AttributeError): pass
        self.reference_artists.clear()
        self.ref_color_index = 0 # カラーサイクルをリセット
        self._update_legend()
        if self.fig and self.fig.canvas:
            self.fig.canvas.draw_idle()

    def _plot_from_list(self):
        if not self.plot_files:
            messagebox.showwarning("リストが空です", "描画するファイルを追加してください。", parent=self)
            return
        if self.canvas_widget: self.canvas_widget.destroy()
        if self.toolbar: self.toolbar.destroy()
        self._clear_reference_artists()
        fig, ax = plt.subplots()
        self.fig = fig
        plot_params = {}
        method = self.method_var.get()
        if method == 'zem3':
            try:
                plots, plot_params = self._process_zem3_for_plotting()
                if not plots: return
                for p in plots: ax.plot(p['x'], p['y'], label=p['label'], marker='o', linestyle='-')
            except Exception as e: messagebox.showerror("ZEM3 処理エラー", f"ZEM3データの処理中にエラーが発生しました。\n\n詳細: {e}", parent=self); return
        else:
            for item in self.plot_files:
                data_path = item['path']
                try:
                    skiprows = self._detect_data_start_row(data_path) if self.auto_detect_header_var.get() else self.skip_rows_var.get() - 1
                    if self.auto_detect_header_var.get(): self.skip_rows_var.set(skiprows + 1)
                    if method == 'iv':
                        df = self._read_iv_data(data_path, skiprows)
                        x_data, y_data, plot_params = self._process_iv_data(df)
                        op, val = item.get('op', '+'), item.get('val', 0.0)
                        y_processed = y_data * val if op == '*' else y_data + val
                        ax.plot(x_data, y_processed, label=data_path.name)
                    elif method in ['xrd', 'raman', 'ftir', 'xps']:
                        plots, plot_params = self._process_spectroscopy_data(data_path, skiprows)
                        for p in plots:
                            op, val = item.get('op', '+'), item.get('val', 0.0)
                            y_processed = p['y'] * val if op == '*' else p['y'] + val
                            ax.plot(p['x'], y_processed, label=p['label'])
                    else: messagebox.showinfo("未対応", f"'{method}' の自動グラフ化はサポートされていません。", parent=self); return
                except Exception as e: messagebox.showerror("処理エラー", f"'{data_path.name}' の処理中にエラーが発生しました。\n\n詳細: {e}", parent=self); continue
        ax.set_xlabel(plot_params.get('xlabel', 'X-axis'))
        ax.set_ylabel(plot_params.get('ylabel', 'Y-axis'))
        ax.grid(True)
        if ax.has_data():
            leg = ax.legend(); leg.set_draggable(True)
        if self.log_scale_var.get(): ax.set_yscale('log')
        try:
            if self.x_min_var.get(): ax.set_xlim(left=float(self.x_min_var.get()))
            if self.x_max_var.get(): ax.set_xlim(right=float(self.x_max_var.get()))
            if self.y_min_var.get(): ax.set_ylim(bottom=float(self.y_min_var.get()))
            if self.y_max_var.get(): ax.set_ylim(top=float(self.y_max_var.get()))
        except ValueError: messagebox.showwarning("入力エラー", "軸範囲には数値を入力してください。", parent=self)
        fig.tight_layout()
        canvas = FigureCanvasTkAgg(fig, master=self.plot_frame)
        self.canvas_widget = canvas.get_tk_widget()
        self.toolbar = NavigationToolbar2Tk(canvas, self.plot_frame)
        self.toolbar.update()
        self.toolbar.pack(side=tk.BOTTOM, fill=tk.X)
        self.canvas_widget.pack(fill=tk.BOTH, expand=True)
        canvas.draw()
        self.canvas_widget.bind("<Control-c>", self._copy_fig_to_clipboard)
        self.canvas_widget.focus_set()

    def _copy_fig_to_clipboard(self, event=None):
        try:
            if sys.platform != 'win32': raise ImportError("Windows以外のOSは現在サポートされていません。")
            import win32clipboard, win32con
            if self.fig is None: messagebox.showwarning("コピー失敗", "コピー対象のグラフがありません。", parent=self); return
            buf = io.BytesIO()
            self.fig.savefig(buf, format='png')
            image = Image.open(buf)
            output_buf = io.BytesIO()
            image.convert('RGB').save(output_buf, 'BMP')
            dib_data = output_buf.getvalue()[14:]
            buf.close(); output_buf.close()
            win32clipboard.OpenClipboard()
            win32clipboard.EmptyClipboard()
            win32clipboard.SetClipboardData(win32con.CF_DIB, dib_data)
            win32clipboard.CloseClipboard()
            messagebox.showinfo("成功", "グラフをクリップボードにコピーしました。", parent=self)
        except ImportError: messagebox.showwarning("機能制限", "クリップボードへの直接コピーはWindowsでのみサポートされています。\nこの機能を使用するには、まずターミナルで `pip install pywin32` を実行してください。\n\nまたは、ツールバーの保存ボタンをご利用ください。", parent=self)
        except Exception as e: messagebox.showwarning("コピー失敗", f"クリップボードへのコピーに失敗しました。\n詳細: {e}", parent=self)

    def _export_to_ppt(self):
        try:
            if self.fig is None: messagebox.showwarning("PPT出力失敗", "出力対象のグラフがありません。", parent=self); return
            image_stream = io.BytesIO()
            self.fig.savefig(image_stream, format='png', dpi=300)
            image_stream.seek(0)
            if self.temp_ppt_path is None: self.temp_ppt_path = Path(os.getenv("TEMP", ".")) / "temp_graph_export.pptx"
            prs = pptx.Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(image_stream, Inches(1), Inches(1), width=Inches(8))
            prs.save(self.temp_ppt_path)
            os.startfile(self.temp_ppt_path)
            messagebox.showinfo("成功", f"PowerPointファイルにグラフを追加しました。\nファイルパス:\n{self.temp_ppt_path}", parent=self)
        except PermissionError: messagebox.showerror("PPT出力失敗", f"PowerPointファイルが使用中です。\nファイルを閉じてから再度お試しください。\n\nファイルパス:\n{self.temp_ppt_path}")
        except Exception as e: messagebox.showerror("PPT出力失敗", f"PowerPointへの出力中にエラーが発生しました。\n詳細: {e}", parent=self)

    def _start_gemini_chat(self):
        if self.fig is None: messagebox.showwarning("Geminiエラー", "チャットを開始するためのグラフがありません。", parent=self); return
        GeminiSettingsDialog(self)
        if not self.gemini_api_key or not self.gemini_model_name: return
        try:
            buf = io.BytesIO(); self.fig.savefig(buf, format='png'); buf.seek(0)
            img = Image.open(buf)
            GeminiChatWindow(self, self.gemini_api_key, self.gemini_model_name, img)
        except Exception as e: messagebox.showerror("エラー", f"チャットウィンドウの準備中にエラーが発生しました:\n{e}", parent=self)

    def _read_iv_data(self, data_path, skiprows):
        suffix = data_path.suffix.lower()
        try:
            if suffix == '.csv': df = pd.read_csv(data_path, header=None, comment='#', on_bad_lines='skip', encoding='cp932', skiprows=skiprows)
            elif suffix in ['.xlsx', '.xls']: df = pd.read_excel(data_path, header=None, skiprows=skiprows)
            elif suffix == '.txt': df = pd.read_csv(data_path, header=None, comment='#', on_bad_lines='skip', encoding='cp932', skiprows=skiprows, sep=r'\s+|,', engine='python')
            else: raise ValueError(f"非対応ファイル形式: {suffix}")
        except Exception as e: raise ValueError(f"ファイルの読み込みに失敗しました。\nエンコーディングや区切り文字を確認してください。\n詳細: {e}")
        if df.shape[1] < 5: raise ValueError(f"IVデータとして読み込むには5列以上必要ですが、{df.shape[1]}列しかありません。\n測定方法の選択が「iv」で正しいか確認してください。")
        df = df.iloc[:, :5]; df.columns = ['Point', 'Time', 'Voltage', 'Current', 'Resistance']
        for col in df.columns[1:]: df[col] = pd.to_numeric(df[col], errors='coerce')
        return df.dropna(subset=['Time', 'Voltage', 'Current'])

    def _process_iv_data(self, df):
        params = {}; x_choice, y_choice = self.x_axis_var.get(), self.y_axis_var.get()
        col_map = {'Time': 'Time', 'Voltage': 'Voltage', 'Current': 'Current', 'Resistance': 'Resistance'}
        unit_map = {'Time': 's', 'Voltage': 'V', 'Current': 'A', 'Resistance': 'Ω'}
        x_data = df[col_map[x_choice]]; params['xlabel'] = f"{x_choice} ({unit_map[x_choice]})"
        if y_choice == 'Current Density':
            area = float(self.area_var.get());
            if area <= 0: raise ValueError("面積は正の数値を入力してください。")
            y_data = df['Current'] / area; params['ylabel'] = "Current Density (A/cm²)"
        else: y_data = df[col_map[y_choice]]; params['ylabel'] = f"{y_choice} ({unit_map[y_choice]})"
        if self.abs_current_var.get() and y_choice in ['Current', 'Current Density']: y_data = y_data.abs(); params['ylabel'] = f"|{params['ylabel']}|"
        return x_data, y_data, params

    def _process_spectroscopy_data(self, data_path, skiprows):
        method = self.method_var.get()
        header_row = skiprows - 1 if skiprows > 0 else None
        try:
            df = pd.read_csv(data_path, comment='#', header=header_row if method == 'xps' else None, skiprows=skiprows if method != 'xps' else None, sep=r'\s+|,|\t', encoding='utf-8', on_bad_lines='skip', engine='python')
        except UnicodeDecodeError:
            df = pd.read_csv(data_path, comment='#', header=header_row if method == 'xps' else None, skiprows=skiprows if method != 'xps' else None, sep=r'\s+|,|\t', encoding='cp932', on_bad_lines='skip', engine='python')
        except Exception as e: raise ValueError(f"ファイルの読み込みに失敗しました。\n詳細: {e}")
        if df.shape[1] < 2: raise ValueError(f"データファイルには少なくとも2列の数値データが必要です。")
        plots, params = [], {}
        if method == 'xps':
            params = {'xlabel': 'Binding Energy (eV)', 'ylabel': 'Intensity (a.u.)'}
            df = df.apply(pd.to_numeric, errors='coerce').dropna()
            x_data = df.iloc[:, 0]
            for col_name in df.columns[1:]:
                plots.append({'x': x_data, 'y': df[col_name], 'label': f'{data_path.name} - {col_name}'})
        else:
            df = df.iloc[:, :2]; df.columns = ['X', 'Y']
            df = df.apply(pd.to_numeric, errors='coerce').dropna()
            x_data, y_data = df['X'], df['Y']
            if method == 'xrd': params = {'xlabel': '2θ (degree)', 'ylabel': 'Intensity (a.u.)'}
            elif method == 'raman': params = {'xlabel': 'Raman Shift (cm-1)', 'ylabel': 'Intensity (a.u.)'}
            elif method == 'ftir':
                params['xlabel'] = 'Wavenumber (cm⁻¹)'
                original_unit, target_unit = self.ftir_original_y_unit_var.get(), self.ftir_target_y_unit_var.get()
                params['ylabel'] = target_unit
                if original_unit != target_unit:
                    if original_unit == 'Absorbance': y_data = 10**(2 - y_data)
                    elif original_unit in ['Transmittance (%)', 'Reflectance (%)']:
                        y_data[y_data <= 0] = 1e-9; y_data = 2 - np.log10(y_data)
            plots.append({'x': x_data, 'y': y_data, 'label': data_path.name})
        return plots, params

    def _get_processed_zem3_df(self):
        all_dfs = []
        for item in self.plot_files:
            try:
                d_actual_m = item['actual_depth_cm'] / 100
                d_file_m = self._read_zem_depth(item['path']) / 100
                if d_file_m <= 0: raise ValueError("ファイルから読み取った深さが0以下です。")
                df_raw = pd.read_csv(item['path'], sep='\t', skiprows=2, header=None, encoding='cp932', engine='python', usecols=[0, 1, 4])
                df_raw.columns = ['Temp', 'Resistivity', 'Seebeck']
                df_raw = df_raw.apply(pd.to_numeric, errors='coerce').dropna()
                df_raw['Resistivity'] *= (d_actual_m / d_file_m)
                df_raw['Conductivity (S/m)'] = 1 / df_raw['Resistivity']
                df_raw['Power Factor (W/mK²)'] = (df_raw['Seebeck'] ** 2) * df_raw['Conductivity (S/m)']
                processed_df = pd.DataFrame({
                    'File Name': item['path'].name, 'Temperature (°C)': df_raw['Temp'],
                    'Electrical conductivity (S/cm)': df_raw['Conductivity (S/m)'] / 100,
                    'Seebeck coefficient (μV/K)': df_raw['Seebeck'] * 1e6,
                    'Power Factor (μW/cmK²)': df_raw['Power Factor (W/mK²)'] * 1e4})
                all_dfs.append(processed_df)
            except Exception as e: messagebox.showerror("処理エラー", f"'{item['path'].name}' の処理中にエラーが発生しました。\n\n詳細: {e}", parent=self); return None
        return pd.concat(all_dfs, ignore_index=True) if all_dfs else None

    def _show_zem3_table_popup(self):
        if not self.plot_files: messagebox.showwarning("リストが空です", "表示するファイルを追加してください。", parent=self); return
        combined_df = self._get_processed_zem3_df()
        if combined_df is not None: ZEM3DataTableWindow(self, combined_df)

    def _process_zem3_for_plotting(self):
        combined_df = self._get_processed_zem3_df()
        if combined_df is None: return [], {}
        plots, x_col, y_col = [], self.x_axis_var.get(), self.y_axis_var.get()
        for file_name, group in combined_df.groupby('File Name'):
            plots.append({'x': group[x_col], 'y': group[y_col], 'label': file_name})
        return plots, {'xlabel': x_col, 'ylabel': y_col}

    def _open_database_search_window(self):
        if not self.fig: messagebox.showwarning("グラフ未描画", "データベース検索の前に、まずグラフを描画してください。", parent=self); return
        method = self.method_var.get()
        if method not in ['xrd', 'raman']: messagebox.showwarning("未対応", f"データベース検索はXRDとRaman測定でのみ利用可能です。", parent=self); return
        
        db_path = self.xrd_db_file if method == 'xrd' else self.raman_db_file
        PeakDatabaseSearchWindow(self, method, db_path)

    def _plot_reference_spectrum(self, ref_df, material_name, add_mode=False):
        """★ 改善版 ★ データベースのスペクトルをリファレンスとしてグラフに描画する"""
        if not add_mode:
            self._clear_reference_artists()
        if not self.fig or not self.fig.axes: return
        ax = self.fig.axes[0]
        
        method = self.method_var.get()
        x_col_name = '2θ (degree)' if method == 'xrd' else 'Raman Shift (cm-1)'
        
        y_min, y_max = ax.get_ylim()
        if y_max <= y_min: y_max = y_min + 1
        
        color = self.ref_colors[self.ref_color_index % len(self.ref_colors)]
        self.ref_color_index += 1
        
        ref_y_scaled = ref_df['Intensity'] / ref_df['Intensity'].max() * (y_max * 0.7)
        
        container = ax.vlines(ref_df[x_col_name], y_min, ref_y_scaled, color=color, linestyle='--', label=f"Ref: {material_name}")
        self.reference_artists.append(container)
        
        for i, row in ref_df.iterrows():
            text = ax.text(row[x_col_name], ref_y_scaled.iloc[i], f" {row['Label']}", color=color, rotation=90, ha='center', va='bottom', size='small')
            self.reference_artists.append(text)
        
        if not add_mode:
            self._update_legend()
            self.fig.canvas.draw_idle()

    def _update_legend(self):
        """★ 新機能 ★ 凡例を再生成して更新する"""
        if not self.fig or not self.fig.axes: return
        ax = self.fig.axes[0]
        if not ax.has_data(): return

        handles, labels = ax.get_legend_handles_labels()
        unique_labels = {}
        for handle, label in zip(handles, labels):
            if label not in unique_labels:
                unique_labels[label] = handle
        
        leg = ax.legend(unique_labels.values(), unique_labels.keys())
        if leg: leg.set_draggable(True)
        
        if self.fig.canvas:
            self.fig.canvas.draw_idle()


# =============================================================================
# === ZEM3データテーブルウィンドウクラス ===
# =============================================================================
class ZEM3DataTableWindow(tb.Toplevel):
    def __init__(self, parent, combined_df):
        super().__init__(parent); self.title("ZEM3 補正後データテーブル & プロッター"); self.geometry("1000x700"); self.combined_df = combined_df
        main_frame = tb.Frame(self, padding=10); main_frame.pack(fill='both', expand=True); main_frame.rowconfigure(0, weight=1); main_frame.columnconfigure(0, weight=1)
        table_frame = tb.LabelFrame(main_frame, text="データテーブル"); table_frame.grid(row=0, column=0, sticky='nsew'); table_frame.rowconfigure(0, weight=1); table_frame.columnconfigure(0, weight=1); self._create_table_view(table_frame)
        stats_frame = tb.LabelFrame(main_frame, text="選択行の統計情報"); stats_frame.grid(row=1, column=0, sticky='ew', pady=(10, 0)); self._create_stats_controls(stats_frame)

    def _create_table_view(self, parent):
        cols = list(self.combined_df.columns); self.tree = tb.Treeview(parent, columns=cols, show='headings', selectmode="extended")
        for col in cols: self.tree.heading(col, text=col); self.tree.column(col, width=200 if "File" in col else 180, anchor='center')
        for _, row in self.combined_df.iterrows(): self.tree.insert("", "end", values=[f"{v:.4g}" if isinstance(v, (float, np.floating)) else v for v in row])
        vsb = tb.Scrollbar(parent, orient="vertical", command=self.tree.yview); vsb.pack(side='right', fill='y'); hsb = tb.Scrollbar(parent, orient="horizontal", command=self.tree.xview); hsb.pack(side='bottom', fill='x')
        self.tree.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set); self.tree.pack(fill='both', expand=True)

    def _create_stats_controls(self, parent):
        self.result_vars = {'Temp': tk.StringVar(value='-'), 'Cond': tk.StringVar(value='-'), 'Seebeck': tk.StringVar(value='-'), 'PF': tk.StringVar(value='-')}
        tk.Button(parent, text="平均・標準偏差を計算", command=self._calculate_stats).grid(row=0, column=0, rowspan=2, padx=10, pady=5, sticky='ns')
        labels = ["Temp (°C)", "Elec. Cond. (S/cm)", "Seebeck Coeff. (μV/K)", "PF (μW/cmK²)"]; var_keys = ["Temp", "Cond", "Seebeck", "PF"]
        for i, label in enumerate(labels): tk.Label(parent, text=f"{label}:").grid(row=0, column=i*2+1, sticky='e', padx=(10, 2)); tk.Label(parent, textvariable=self.result_vars[var_keys[i]]).grid(row=0, column=i*2+2, sticky='w')

    def _calculate_stats(self):
        selected_items = self.tree.selection()
        if not selected_items: messagebox.showwarning("行未選択", "統計を計算する行を1つ以上選択してください。", parent=self); return
        col_map = {'Temperature (°C)': 'Temp', 'Electrical conductivity (S/cm)': 'Cond', 'Seebeck coefficient (μV/K)': 'Seebeck', 'Power Factor (μW/cmK²)': 'PF'}
        selected_data = defaultdict(list)
        for item in selected_items:
            row = self.tree.item(item, 'values')
            for idx, col_name in enumerate(self.combined_df.columns):
                if col_name in col_map:
                    try: selected_data[col_name].append(float(row[idx]))
                    except (ValueError, IndexError): pass
        for col_name, key in col_map.items():
            data = selected_data.get(col_name, [])
            if data:
                avg, std = np.mean(data), np.std(data, ddof=1) if len(data) > 1 else 0
                self.result_vars[key].set(f"{avg:.4g} ± {std:.4g}")
            else: self.result_vars[key].set("-")

# =============================================================================
# === アプリケーションの実行 ===
# =============================================================================
if __name__ == "__main__":
    app = DataAnalyzerApp()
    app.mainloop()
